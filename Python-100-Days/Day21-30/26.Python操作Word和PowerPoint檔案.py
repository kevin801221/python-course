#!/usr/bin/env python3
"""
從 26.Python操作Word和PowerPoint文件.md 提取的 Python 範例代碼
"""

# === 範例 1 ===
from docx import Document
from docx.shared import Cm, Pt

from docx.document import Document as Doc

# 建立代表Word文件的Doc物件
document = Document()  # type: Doc
# 新增大標題
document.add_heading('快快樂樂學Python', 0)
# 新增段落
p = document.add_paragraph('Python是一門非常流行的程式語言，它')
run = p.add_run('簡單')
run.bold = True
run.font.size = Pt(18)
p.add_run('而且')
run = p.add_run('優雅')
run.font.size = Pt(18)
run.underline = True
p.add_run('。')

# 新增一級標題
document.add_heading('Heading, level 1', level=1)
# 新增帶樣式的段落
document.add_paragraph('Intense quote', style='Intense Quote')
# 新增無序列表
document.add_paragraph(
    'first item in unordered list', style='List Bullet'
)
document.add_paragraph(
    'second item in ordered list', style='List Bullet'
)
# 新增有序列表
document.add_paragraph(
    'first item in ordered list', style='List Number'
)
document.add_paragraph(
    'second item in ordered list', style='List Number'
)

# 新增圖片（注意路徑和圖片必須要存在）
document.add_picture('resources/guido.jpg', width=Cm(5.2))

# 新增分節符
document.add_section()

records = (
    ('Kevin', '男', '1995-5-5'),
    ('孫美麗', '女', '1992-2-2')
)
# 新增表格
table = document.add_table(rows=1, cols=3)
table.style = 'Dark List'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = '姓名'
hdr_cells[1].text = '性別'
hdr_cells[2].text = '出生日期'
# 為表格新增行
for name, sex, birthday in records:
    row_cells = table.add_row().cells
    row_cells[0].text = name
    row_cells[1].text = sex
    row_cells[2].text = birthday

# 新增分頁符
document.add_page_break()

# 儲存文件
document.save('demo.docx')
# === 範例 2 ===
from docx import Document
from docx.document import Document as Doc

doc = Document('resources/離職證明.docx')  # type: Doc
for no, p in enumerate(doc.paragraphs):
    print(no, p.text)
# === 範例 3 ===
from docx import Document
from docx.document import Document as Doc

# 將真實資訊用字典的方式儲存在列表中
employees = [
    {
        'name': 'Kevin',
        'id': '100200198011280001',
        'sdate': '2008年3月1日',
        'edate': '2012年2月29日',
        'department': '產品研發',
        'position': '架構師',
        'company': '成都華為技術有限公司'
    },
    {
        'name': '王大錘',
        'id': '510210199012125566',
        'sdate': '2019年1月1日',
        'edate': '2021年4月30日',
        'department': '產品研發',
        'position': 'Python開發工程師',
        'company': '成都穀道科技有限公司'
    },
    {
        'name': '李元芳',
        'id': '2102101995103221599',
        'sdate': '2020年5月10日',
        'edate': '2021年3月5日',
        'department': '產品研發',
        'position': 'Java開發工程師',
        'company': '同城企業管理集團有限公司'
    },
]
# 對列表進行迴圈遍歷，批次生成Word文件 
for emp_dict in employees:
    # 讀取離職證明模板檔案
    doc = Document('resources/離職證明模板.docx')  # type: Doc
    # 迴圈遍歷所有段落尋找佔位符
    for p in doc.paragraphs:
        if '{' not in p.text:
            continue
        # 不能直接修改段落內容，否則會丟失樣式
        # 所以需要對段落中的元素進行遍歷並進行查詢替換
        for run in p.runs:
            if '{' not in run.text:
                continue
            # 將佔位符換成實際內容
            start, end = run.text.find('{'), run.text.find('}')
            key, place_holder = run.text[start + 1:end], run.text[start:end + 1]
            run.text = run.text.replace(place_holder, emp_dict[key])
    # 每個人對應儲存一個Word文件
    doc.save(f'{emp_dict["name"]}離職證明.docx')
# === 範例 4 ===
from pptx import Presentation

# 建立幻燈片物件
pres = Presentation()

# 選擇母版新增一頁
title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)
# 獲取標題欄和副標題欄
title = slide.shapes.title
subtitle = slide.placeholders[1]
# 編輯標題和副標題
title.text = "Welcome to Python"
subtitle.text = "Life is short, I use Python"

# 選擇母版新增一頁
bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)
# 獲取頁面上所有形狀
shapes = slide.shapes
# 獲取標題和主體
title_shape = shapes.title
body_shape = shapes.placeholders[1]
# 編輯標題
title_shape.text = 'Introduction'
# 編輯主體內容
tf = body_shape.text_frame
tf.text = 'History of Python'
# 新增一個一級段落
p = tf.add_paragraph()
p.text = 'X\'max 1989'
p.level = 1
# 新增一個二級段落
p = tf.add_paragraph()
p.text = 'Guido began to write interpreter for Python.'
p.level = 2

# 儲存幻燈片
pres.save('test.pptx')
